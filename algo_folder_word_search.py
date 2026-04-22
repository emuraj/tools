# document_keyword_search_gui.py
"""
Purpose:
    Desktop tkinter application for scanning large folders of supported document
    types and counting exact whole-word keyword hits.

What this file does:
    - Lets the user choose an input folder and an output folder
    - Accepts slash-delimited keywords, such as:
          maintenance/calibration/cleaning
    - Lets the user choose file types by checkbox:
          .docx, .doc, .pdf, .xlsx, .ppt, .pptx
    - Recursively scans only the selected file types
    - Counts case-insensitive exact whole-word matches
    - Treats punctuation around a word as acceptable
    - Does not count longer words or hyphenated compounds as exact matches
    - Writes results to a CSV file
    - Writes a formal PDF run log with:
          preamble
          purpose
          scope
          executive summary
          machine stats
          run stats
          detailed event log
          headers / footers / page X of Y

Place in the larger scheme:
    This is a standalone desktop utility intended for fast, audit-friendly bulk
    keyword searching of mixed document populations with operator-friendly GUI
    controls and exportable traceability outputs.

Why that matters:
    The tool is meant to support structured document review at scale while
    preserving a clean execution narrative and usable outputs for downstream
    sorting, triage, and inspection-style review.
"""

from __future__ import annotations

import csv
import getpass
import os
import platform
import queue
import re
import socket
import statistics
import sys
import tempfile
import threading
import time
import traceback
from dataclasses import dataclass
from datetime import datetime
from pathlib import Path
from typing import Callable, Dict, List, Optional

import tkinter as tk
from tkinter import filedialog, messagebox, scrolledtext, ttk

from docx import Document
from openpyxl import load_workbook
from pypdf import PdfReader
from pptx import Presentation
from reportlab.lib import colors
from reportlab.lib.pagesizes import letter
from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
from reportlab.lib.units import inch
from reportlab.pdfbase.pdfmetrics import stringWidth
from reportlab.platypus import Paragraph, SimpleDocTemplate, Spacer, Table, TableStyle
from reportlab.pdfgen import canvas as pdf_canvas


APP_TITLE = "Document Whole-Word Keyword Search Tool"
WINDOW_SIZE = "1180x820"
CSV_PREFIX = "algo_folder_word_search"
PDF_PREFIX = "algo_folder_word_search_log"
TOOL_VERSION = "1.0"
SUPPORTED_TYPES = [".docx", ".doc", ".pdf", ".xlsx", ".ppt", ".pptx"]


@dataclass
class AppPaths:
    input_folder: Path
    output_folder: Path


@dataclass
class RunStats:
    run_id: str
    timestamp_compact: str
    timestamp_human: str
    start_dt: datetime
    end_dt: Optional[datetime]
    elapsed_s: float
    total_discovered: int
    processed: int
    ok: int
    failed: int
    docs_per_min: float
    total_keyword_hits: int
    keywords: List[str]
    file_types_selected: List[str]
    input_folder: str
    output_folder: str
    host_name: str
    user_name: str
    operating_system: str
    python_version: str
    processor_count: int
    machine_arch: str
    per_type_counts: Dict[str, int]
    prefix_summary: Dict[str, int]
    hits_stats: Dict[str, float]


class NumberedCanvas(pdf_canvas.Canvas):
    """
    ReportLab canvas that supports Page X of Y and consistent header/footer.
    """

    def __init__(self, *args, **kwargs):
        self.report_title = kwargs.pop("report_title", "")
        self.run_id = kwargs.pop("run_id", "")
        self.generated_ts = kwargs.pop("generated_ts", "")
        self.file_label = kwargs.pop("file_label", "")
        super().__init__(*args, **kwargs)
        self._saved_page_states = []

    def showPage(self):
        self._saved_page_states.append(dict(self.__dict__))
        self._startPage()

    def save(self):
        page_count = len(self._saved_page_states)
        for state in self._saved_page_states:
            self.__dict__.update(state)
            self._draw_header_footer(page_count)
            super().showPage()
        super().save()

    def _draw_header_footer(self, page_count: int):
        page_num = self._pageNumber
        width, height = letter

        header_y = height - 36
        footer_y = 24

        self.setFont("Helvetica-Bold", 9)
        self.drawString(36, header_y, self.report_title)

        run_id_text = f"Run ID: {self.run_id}"
        run_id_width = stringWidth(run_id_text, "Helvetica-Bold", 9)
        self.drawString(width - 36 - run_id_width, header_y, run_id_text)

        self.setStrokeColor(colors.grey)
        self.setLineWidth(0.5)
        self.line(36, header_y - 6, width - 36, header_y - 6)
        self.line(36, footer_y + 10, width - 36, footer_y + 10)

        self.setFont("Helvetica", 8)
        self.drawString(36, footer_y, f"Generated: {self.generated_ts}")

        center_text = self.file_label
        center_width = stringWidth(center_text, "Helvetica", 8)
        self.drawString((width - center_width) / 2, footer_y, center_text)

        page_text = f"Page {page_num} of {page_count}"
        page_width = stringWidth(page_text, "Helvetica", 8)
        self.drawString(width - 36 - page_width, footer_y, page_text)


def make_timestamp_compact() -> str:
    return datetime.now().strftime("%m%d%Y%H%M%S")


def make_timestamp_human(dt: Optional[datetime] = None) -> str:
    actual = dt or datetime.now()
    return actual.strftime("%d%b%Y %H:%M:%S").upper()


def make_run_id(dt: Optional[datetime] = None) -> str:
    actual = dt or datetime.now()
    return f"DWS-{actual.strftime('%Y%m%d-%H%M%S')}"


def timestamped_log_message(message: str) -> str:
    return f"[{datetime.now().strftime('%H:%M:%S')}] {message}"


def parse_keywords(raw_keywords: str) -> List[str]:
    parts = [part.strip() for part in raw_keywords.split("/") if part.strip()]
    seen = set()
    ordered: List[str] = []
    for part in parts:
        norm = part.casefold()
        if norm not in seen:
            seen.add(norm)
            ordered.append(part)
    return ordered


def build_keyword_patterns(keywords: List[str]) -> Dict[str, re.Pattern]:
    patterns: Dict[str, re.Pattern] = {}
    for keyword in keywords:
        escaped = re.escape(keyword)
        patterns[keyword] = re.compile(rf"(?<![\w-]){escaped}(?![\w-])", re.IGNORECASE)
    return patterns


def extract_prefix(file_name: str) -> str:
    stem = Path(file_name).stem.strip()

    if "-" in stem:
        return stem.split("-", 1)[0].upper()

    if re.match(r"^F\d+[A-Z]*$", stem, re.IGNORECASE):
        return "F"

    leading_letters = re.match(r"^([A-Za-z]+)", stem)
    if leading_letters:
        return leading_letters.group(1).upper()

    return "UNKNOWN"


def compute_total_hits(row: Dict[str, int | str], keywords: List[str]) -> int:
    return sum(int(row.get(keyword, 0)) for keyword in keywords)


def sort_rows_by_prefix_and_total(rows: List[Dict[str, int | str]]) -> List[Dict[str, int | str]]:
    return sorted(
        rows,
        key=lambda row: (
            str(row["Prefix"]).upper(),
            -int(row["Total Hits"]),
            str(row["File Name"]).upper(),
        ),
    )


def extract_docx_text(path: Path) -> str:
    doc = Document(str(path))
    pieces: List[str] = []

    for paragraph in doc.paragraphs:
        if paragraph.text:
            pieces.append(paragraph.text)

    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                for paragraph in cell.paragraphs:
                    if paragraph.text:
                        pieces.append(paragraph.text)

    return "\n".join(pieces)


def extract_pdf_text(path: Path) -> str:
    reader = PdfReader(str(path))
    pieces: List[str] = []
    for page in reader.pages:
        text = page.extract_text()
        if text:
            pieces.append(text)
    return "\n".join(pieces)


def extract_xlsx_text(path: Path) -> str:
    wb = load_workbook(filename=str(path), read_only=True, data_only=True)
    pieces: List[str] = []
    for sheet in wb.worksheets:
        pieces.append(f"[Worksheet: {sheet.title}]")
        for row in sheet.iter_rows(values_only=True):
            row_values = [str(cell) for cell in row if cell is not None and str(cell).strip()]
            if row_values:
                pieces.append(" | ".join(row_values))
    wb.close()
    return "\n".join(pieces)


def extract_pptx_text(path: Path) -> str:
    prs = Presentation(str(path))
    pieces: List[str] = []
    for slide_index, slide in enumerate(prs.slides, start=1):
        pieces.append(f"[Slide {slide_index}]")
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                pieces.append(shape.text)
    return "\n".join(pieces)


def extract_doc_text_via_word_com(path: Path) -> str:
    """
    Extract text from legacy .doc using Microsoft Word COM automation.
    Requires:
        - Windows
        - Microsoft Word installed
        - pywin32 installed
    """
    try:
        import pythoncom
        import win32com.client  # type: ignore
        from win32com.client import constants  # type: ignore
    except Exception as exc:
        raise RuntimeError(
            "Legacy .doc support requires Windows, Microsoft Word, and pywin32."
        ) from exc

    pythoncom.CoInitialize()
    word = None
    doc = None
    temp_docx = None

    try:
        word = win32com.client.DispatchEx("Word.Application")
        word.Visible = False
        word.DisplayAlerts = 0

        doc = word.Documents.Open(str(path), ReadOnly=True)
        temp_dir = Path(tempfile.mkdtemp(prefix="doc_extract_"))
        temp_docx = temp_dir / f"{path.stem}_temp_extract.docx"

        wdFormatXMLDocument = 12
        doc.SaveAs(str(temp_docx), FileFormat=wdFormatXMLDocument)
        doc.Close(False)
        doc = None

        return extract_docx_text(temp_docx)

    finally:
        try:
            if doc is not None:
                doc.Close(False)
        except Exception:
            pass
        try:
            if word is not None:
                word.Quit()
        except Exception:
            pass
        try:
            pythoncom.CoUninitialize()
        except Exception:
            pass
        try:
            if temp_docx and temp_docx.exists():
                temp_docx.unlink(missing_ok=True)
                temp_docx.parent.rmdir()
        except Exception:
            pass


def extract_ppt_text_via_powerpoint_com(path: Path) -> str:
    """
    Extract text from legacy .ppt using PowerPoint COM automation.
    Requires:
        - Windows
        - Microsoft PowerPoint installed
        - pywin32 installed
    """
    try:
        import pythoncom
        import win32com.client  # type: ignore
    except Exception as exc:
        raise RuntimeError(
            "Legacy .ppt support requires Windows, Microsoft PowerPoint, and pywin32."
        ) from exc

    pythoncom.CoInitialize()
    app = None
    presentation = None
    temp_pptx = None

    try:
        app = win32com.client.DispatchEx("PowerPoint.Application")
        app.Visible = 1

        presentation = app.Presentations.Open(str(path), WithWindow=False)
        temp_dir = Path(tempfile.mkdtemp(prefix="ppt_extract_"))
        temp_pptx = temp_dir / f"{path.stem}_temp_extract.pptx"

        ppSaveAsOpenXMLPresentation = 24
        presentation.SaveAs(str(temp_pptx), ppSaveAsOpenXMLPresentation)
        presentation.Close()
        presentation = None

        return extract_pptx_text(temp_pptx)

    finally:
        try:
            if presentation is not None:
                presentation.Close()
        except Exception:
            pass
        try:
            if app is not None:
                app.Quit()
        except Exception:
            pass
        try:
            pythoncom.CoUninitialize()
        except Exception:
            pass
        try:
            if temp_pptx and temp_pptx.exists():
                temp_pptx.unlink(missing_ok=True)
                temp_pptx.parent.rmdir()
        except Exception:
            pass


def get_extractor_for_extension(ext: str) -> Callable[[Path], str]:
    ext_lower = ext.lower()
    if ext_lower == ".docx":
        return extract_docx_text
    if ext_lower == ".doc":
        return extract_doc_text_via_word_com
    if ext_lower == ".pdf":
        return extract_pdf_text
    if ext_lower == ".xlsx":
        return extract_xlsx_text
    if ext_lower == ".pptx":
        return extract_pptx_text
    if ext_lower == ".ppt":
        return extract_ppt_text_via_powerpoint_com
    raise ValueError(f"Unsupported file type: {ext}")


def find_matching_files(root_folder: Path, selected_types: List[str]) -> List[Path]:
    selected = {ext.lower() for ext in selected_types}
    all_matches: List[Path] = []
    for path in root_folder.rglob("*"):
        if path.is_file() and path.suffix.lower() in selected and not path.name.startswith("~$"):
            all_matches.append(path)
    return sorted(all_matches, key=lambda p: str(p).lower())


def write_csv(output_csv_path: Path, keywords: List[str], rows: List[Dict[str, int | str]]) -> None:
    fieldnames = ["File Name", "Prefix", "File Type", "Total Hits", *keywords]
    with output_csv_path.open("w", newline="", encoding="utf-8-sig") as csvfile:
        writer = csv.DictWriter(csvfile, fieldnames=fieldnames)
        writer.writeheader()
        for row in rows:
            writer.writerow(row)


def safe_float(value: float) -> str:
    return f"{value:.2f}"


def build_hits_stats(rows: List[Dict[str, int | str]]) -> Dict[str, float]:
    totals = [int(row["Total Hits"]) for row in rows]
    if not totals:
        return {
            "min": 0.0,
            "median": 0.0,
            "max": 0.0,
            "mean": 0.0,
        }
    return {
        "min": float(min(totals)),
        "median": float(statistics.median(totals)),
        "max": float(max(totals)),
        "mean": float(statistics.mean(totals)),
    }


def build_pdf_report(
    output_pdf_path: Path,
    stats: RunStats,
    log_lines: List[str],
) -> None:
    styles = getSampleStyleSheet()

    title_style = ParagraphStyle(
        "TitleCustom",
        parent=styles["Title"],
        fontName="Helvetica-Bold",
        fontSize=16,
        leading=20,
        spaceAfter=12,
    )

    heading_style = ParagraphStyle(
        "HeadingCustom",
        parent=styles["Heading2"],
        fontName="Helvetica-Bold",
        fontSize=11,
        leading=14,
        spaceBefore=10,
        spaceAfter=6,
    )

    body_style = ParagraphStyle(
        "BodyCustom",
        parent=styles["BodyText"],
        fontName="Helvetica",
        fontSize=9,
        leading=12,
        spaceAfter=6,
    )

    mono_style = ParagraphStyle(
        "MonoCustom",
        parent=styles["BodyText"],
        fontName="Courier",
        fontSize=8,
        leading=10,
        spaceAfter=2,
    )

    elements: List = []

    elements.append(Paragraph("Document Whole-Word Keyword Search Run Report", title_style))
    elements.append(Paragraph(f"Run ID: {stats.run_id}", body_style))
    elements.append(Paragraph(f"Run Timestamp: {stats.timestamp_human}", body_style))
    elements.append(Spacer(1, 8))

    purpose = (
        "This report documents execution of the Document Whole-Word Keyword Search Tool, "
        "a desktop utility designed to scan a selected source folder for user-selected "
        "document file types, search those files for user-specified exact whole-word "
        "keywords, and export structured hit counts by file to CSV."
    )
    scope = (
        "This run evaluated all qualifying files discovered recursively beneath the "
        "selected input folder for the file types checked by the operator. Matching was "
        "performed case-insensitively using exact whole-word logic. Punctuation adjacent "
        "to a word was treated as a valid hit boundary. Hyphenated compounds and longer "
        "word variants were excluded from hit counts."
    )
    output_intent = (
        "The CSV output is intended to support document-family review by grouping records "
        "according to filename prefix and sorting files within each prefix by descending "
        "total keyword hits. This PDF captures the execution narrative, runtime environment, "
        "processing statistics, warnings, and final output locations for traceability."
    )
    executive_summary = (
        f"This run processed {stats.total_discovered} discovered file(s) across the selected "
        f"file types and completed in {safe_float(stats.elapsed_s)} seconds. A total of "
        f"{stats.total_keyword_hits} keyword hit(s) were identified across {stats.ok} "
        f"successfully processed file(s). Output files were generated to the designated "
        f"output directory, including a CSV hit table and this PDF run log. "
        f"{stats.failed} file(s) could not be read and were recorded as warnings."
    )

    elements.append(Paragraph("Purpose", heading_style))
    elements.append(Paragraph(purpose, body_style))
    elements.append(Paragraph("Scope", heading_style))
    elements.append(Paragraph(scope, body_style))
    elements.append(Paragraph("Output Intent", heading_style))
    elements.append(Paragraph(output_intent, body_style))
    elements.append(Paragraph("Executive Summary", heading_style))
    elements.append(Paragraph(executive_summary, body_style))

    elements.append(Paragraph("Runtime Environment", heading_style))
    runtime_table = Table(
        [
            ["Host Name", stats.host_name],
            ["User Name", stats.user_name],
            ["Operating System", stats.operating_system],
            ["Python Version", stats.python_version],
            ["Processor Count", str(stats.processor_count)],
            ["Machine Architecture", stats.machine_arch],
            ["Tool Version", TOOL_VERSION],
        ],
        colWidths=[1.8 * inch, 4.8 * inch],
    )
    runtime_table.setStyle(
        TableStyle(
            [
                ("GRID", (0, 0), (-1, -1), 0.3, colors.grey),
                ("BACKGROUND", (0, 0), (0, -1), colors.whitesmoke),
                ("FONTNAME", (0, 0), (-1, -1), "Helvetica"),
                ("FONTSIZE", (0, 0), (-1, -1), 8),
                ("VALIGN", (0, 0), (-1, -1), "TOP"),
            ]
        )
    )
    elements.append(runtime_table)
    elements.append(Spacer(1, 8))

    elements.append(Paragraph("Run Statistics", heading_style))
    run_stats_table = Table(
        [
            ["Start Time", make_timestamp_human(stats.start_dt)],
            ["End Time", make_timestamp_human(stats.end_dt)],
            ["Elapsed Seconds", safe_float(stats.elapsed_s)],
            ["Total Discovered", str(stats.total_discovered)],
            ["Processed", str(stats.processed)],
            ["Successful", str(stats.ok)],
            ["Failed", str(stats.failed)],
            ["Files Per Minute", safe_float(stats.docs_per_min)],
            ["Total Keyword Hits", str(stats.total_keyword_hits)],
            ["Keywords Used", " / ".join(stats.keywords)],
            ["File Types Selected", " / ".join(stats.file_types_selected)],
            ["Input Folder", stats.input_folder],
            ["Output Folder", stats.output_folder],
            [
                "Per-File Total Hit Stats",
                f"min={safe_float(stats.hits_stats['min'])}, "
                f"median={safe_float(stats.hits_stats['median'])}, "
                f"mean={safe_float(stats.hits_stats['mean'])}, "
                f"max={safe_float(stats.hits_stats['max'])}",
            ],
        ],
        colWidths=[1.8 * inch, 4.8 * inch],
    )
    run_stats_table.setStyle(
        TableStyle(
            [
                ("GRID", (0, 0), (-1, -1), 0.3, colors.grey),
                ("BACKGROUND", (0, 0), (0, -1), colors.whitesmoke),
                ("FONTNAME", (0, 0), (-1, -1), "Helvetica"),
                ("FONTSIZE", (0, 0), (-1, -1), 8),
                ("VALIGN", (0, 0), (-1, -1), "TOP"),
            ]
        )
    )
    elements.append(run_stats_table)
    elements.append(Spacer(1, 8))

    elements.append(Paragraph("Discovered by File Type", heading_style))
    per_type_rows = [["File Type", "Count"]]
    for ext in sorted(stats.per_type_counts):
        per_type_rows.append([ext, str(stats.per_type_counts[ext])])
    per_type_table = Table(per_type_rows, colWidths=[2.0 * inch, 1.2 * inch])
    per_type_table.setStyle(
        TableStyle(
            [
                ("GRID", (0, 0), (-1, -1), 0.3, colors.grey),
                ("BACKGROUND", (0, 0), (-1, 0), colors.lightgrey),
                ("FONTNAME", (0, 0), (-1, -1), "Helvetica"),
                ("FONTSIZE", (0, 0), (-1, -1), 8),
            ]
        )
    )
    elements.append(per_type_table)
    elements.append(Spacer(1, 8))

    elements.append(Paragraph("Prefix Total Summary", heading_style))
    prefix_rows = [["Prefix", "Total Keyword Hits"]]
    for prefix in sorted(stats.prefix_summary):
        prefix_rows.append([prefix, str(stats.prefix_summary[prefix])])
    prefix_table = Table(prefix_rows, colWidths=[2.0 * inch, 1.5 * inch])
    prefix_table.setStyle(
        TableStyle(
            [
                ("GRID", (0, 0), (-1, -1), 0.3, colors.grey),
                ("BACKGROUND", (0, 0), (-1, 0), colors.lightgrey),
                ("FONTNAME", (0, 0), (-1, -1), "Helvetica"),
                ("FONTSIZE", (0, 0), (-1, -1), 8),
            ]
        )
    )
    elements.append(prefix_table)
    elements.append(Spacer(1, 8))

    elements.append(Paragraph("Detailed Event Log", heading_style))
    for line in log_lines:
        safe_line = line.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")
        elements.append(Paragraph(safe_line, mono_style))

    doc = SimpleDocTemplate(
        str(output_pdf_path),
        pagesize=letter,
        leftMargin=36,
        rightMargin=36,
        topMargin=54,
        bottomMargin=40,
    )

    doc.build(
        elements,
        canvasmaker=lambda *args, **kwargs: NumberedCanvas(
            *args,
            report_title="Document Whole-Word Keyword Search Run Report",
            run_id=stats.run_id,
            generated_ts=stats.timestamp_human,
            file_label=output_pdf_path.name,
            **kwargs,
        ),
    )


class DocumentKeywordSearchApp:
    def __init__(self, root: tk.Tk) -> None:
        self.root = root
        self.root.title(APP_TITLE)
        self.root.geometry(WINDOW_SIZE)
        self.root.minsize(1040, 740)

        self.worker_thread: Optional[threading.Thread] = None
        self.message_queue: queue.Queue = queue.Queue()
        self.log_lines: List[str] = []
        self.is_running = False

        self.input_var = tk.StringVar()
        self.output_var = tk.StringVar()
        self.keywords_var = tk.StringVar()

        self.file_type_vars: Dict[str, tk.BooleanVar] = {
            ".docx": tk.BooleanVar(value=True),
            ".doc": tk.BooleanVar(value=False),
            ".pdf": tk.BooleanVar(value=False),
            ".xlsx": tk.BooleanVar(value=False),
            ".ppt": tk.BooleanVar(value=False),
            ".pptx": tk.BooleanVar(value=False),
        }

        self._build_gui()
        self._log("Ready.")
        self._poll_queue()

    def _build_gui(self) -> None:
        self.root.grid_rowconfigure(1, weight=1)
        self.root.grid_columnconfigure(0, weight=1)

        top_frame = ttk.Frame(self.root, padding=12)
        top_frame.grid(row=0, column=0, sticky="nsew")
        top_frame.grid_columnconfigure(1, weight=1)

        ttk.Label(top_frame, text="Input Folder:").grid(row=0, column=0, sticky="w", padx=(0, 8), pady=6)
        self.input_entry = ttk.Entry(top_frame, textvariable=self.input_var)
        self.input_entry.grid(row=0, column=1, sticky="ew", pady=6)
        ttk.Button(top_frame, text="Browse", command=self._browse_input_folder).grid(row=0, column=2, sticky="ew", pady=6)

        ttk.Label(top_frame, text="Output Folder:").grid(row=1, column=0, sticky="w", padx=(0, 8), pady=6)
        self.output_entry = ttk.Entry(top_frame, textvariable=self.output_var)
        self.output_entry.grid(row=1, column=1, sticky="ew", pady=6)
        ttk.Button(top_frame, text="Browse", command=self._browse_output_folder).grid(row=1, column=2, sticky="ew", pady=6)

        ttk.Label(top_frame, text="Keywords (/ separated):").grid(row=2, column=0, sticky="w", padx=(0, 8), pady=6)
        self.keywords_entry = ttk.Entry(top_frame, textvariable=self.keywords_var)
        self.keywords_entry.grid(row=2, column=1, columnspan=2, sticky="ew", pady=6)

        file_types_frame = ttk.LabelFrame(top_frame, text="File Types to Include", padding=8)
        file_types_frame.grid(row=3, column=0, columnspan=3, sticky="ew", pady=(8, 8))

        col = 0
        for ext in SUPPORTED_TYPES:
            chk = ttk.Checkbutton(file_types_frame, text=ext, variable=self.file_type_vars[ext])
            chk.grid(row=0, column=col, padx=(0, 14), sticky="w")
            col += 1

        rules_text = (
            "Matching Rules: Exact whole word, case-insensitive, punctuation allowed around the word, "
            "hyphenated compounds excluded"
        )
        ttk.Label(top_frame, text=rules_text).grid(row=4, column=0, columnspan=3, sticky="w", pady=(2, 10))

        button_frame = ttk.Frame(top_frame)
        button_frame.grid(row=5, column=0, columnspan=3, sticky="w")

        self.run_button = ttk.Button(button_frame, text="Run Search", command=self._start_search)
        self.run_button.grid(row=0, column=0, padx=(0, 8))

        self.clear_log_button = ttk.Button(button_frame, text="Clear Log", command=self._clear_log)
        self.clear_log_button.grid(row=0, column=1, padx=(0, 8))

        self.progress_label = ttk.Label(button_frame, text="Idle")
        self.progress_label.grid(row=0, column=2, padx=(12, 0))

        self.progress_bar = ttk.Progressbar(top_frame, mode="determinate")
        self.progress_bar.grid(row=6, column=0, columnspan=3, sticky="ew", pady=(10, 0))

        log_frame = ttk.LabelFrame(self.root, text="Event Log", padding=10)
        log_frame.grid(row=1, column=0, sticky="nsew", padx=12, pady=(0, 12))
        log_frame.grid_rowconfigure(0, weight=1)
        log_frame.grid_columnconfigure(0, weight=1)

        self.log_widget = scrolledtext.ScrolledText(
            log_frame,
            wrap=tk.WORD,
            state=tk.DISABLED,
            font=("Consolas", 10),
        )
        self.log_widget.grid(row=0, column=0, sticky="nsew")

    def _browse_input_folder(self) -> None:
        folder = filedialog.askdirectory(title="Select Input Folder")
        if folder:
            self.input_var.set(folder)
            self._log(f"Input folder selected: {folder}")

    def _browse_output_folder(self) -> None:
        folder = filedialog.askdirectory(title="Select Output Folder")
        if folder:
            self.output_var.set(folder)
            self._log(f"Output folder selected: {folder}")

    def _clear_log(self) -> None:
        if self.is_running:
            messagebox.showwarning("Search Running", "Cannot clear log while a search is running.")
            return
        self.log_lines.clear()
        self.log_widget.configure(state=tk.NORMAL)
        self.log_widget.delete("1.0", tk.END)
        self.log_widget.configure(state=tk.DISABLED)
        self._log("Log cleared. Ready.")

    def _get_selected_file_types(self) -> List[str]:
        return [ext for ext, var in self.file_type_vars.items() if var.get()]

    def _set_running_state(self, running: bool) -> None:
        self.is_running = running
        state = tk.DISABLED if running else tk.NORMAL

        self.run_button.configure(state=state)
        self.clear_log_button.configure(state=state)
        self.input_entry.configure(state=state)
        self.output_entry.configure(state=state)
        self.keywords_entry.configure(state=state)

    def _validate_inputs(self) -> Optional[tuple[AppPaths, List[str], List[str]]]:
        input_text = self.input_var.get().strip()
        output_text = self.output_var.get().strip()
        keyword_text = self.keywords_var.get().strip()
        selected_types = self._get_selected_file_types()

        if not input_text:
            messagebox.showerror("Missing Input Folder", "Please select an input folder.")
            return None

        if not output_text:
            messagebox.showerror("Missing Output Folder", "Please select an output folder.")
            return None

        if not keyword_text:
            messagebox.showerror("Missing Keywords", "Please enter one or more keywords separated by '/'.")
            return None

        if not selected_types:
            messagebox.showerror("Missing File Types", "Please select at least one file type.")
            return None

        input_folder = Path(input_text)
        output_folder = Path(output_text)

        if not input_folder.exists() or not input_folder.is_dir():
            messagebox.showerror("Invalid Input Folder", "The selected input folder does not exist or is not a folder.")
            return None

        if not output_folder.exists() or not output_folder.is_dir():
            messagebox.showerror("Invalid Output Folder", "The selected output folder does not exist or is not a folder.")
            return None

        keywords = parse_keywords(keyword_text)
        if not keywords:
            messagebox.showerror("Invalid Keywords", "No valid keywords were found.")
            return None

        return AppPaths(input_folder=input_folder, output_folder=output_folder), keywords, selected_types

    def _start_search(self) -> None:
        validated = self._validate_inputs()
        if validated is None:
            return

        paths, keywords, selected_types = validated

        self._set_running_state(True)
        self.progress_label.configure(text="Running")
        self.progress_bar["value"] = 0
        self.progress_bar["maximum"] = 100

        self._log(f"Keywords loaded: {', '.join(keywords)}")
        self._log(f"File types selected: {', '.join(selected_types)}")
        self._log("Starting scan...")

        self.worker_thread = threading.Thread(
            target=self._run_search_worker,
            args=(paths, keywords, selected_types),
            daemon=True,
        )
        self.worker_thread.start()

    def _run_search_worker(self, paths: AppPaths, keywords: List[str], selected_types: List[str]) -> None:
        start_dt = datetime.now()
        start_perf = time.perf_counter()

        timestamp_compact = start_dt.strftime("%m%d%Y%H%M%S")
        timestamp_human = make_timestamp_human(start_dt)
        run_id = make_run_id(start_dt)

        output_csv = paths.output_folder / f"{CSV_PREFIX}_{timestamp_compact}.csv"
        output_pdf = paths.output_folder / f"{PDF_PREFIX}_{timestamp_compact}.pdf"

        try:
            patterns = build_keyword_patterns(keywords)
            matching_files = find_matching_files(paths.input_folder, selected_types)
            total_files = len(matching_files)

            per_type_counts: Dict[str, int] = {ext: 0 for ext in selected_types}
            for path in matching_files:
                ext = path.suffix.lower()
                per_type_counts[ext] = per_type_counts.get(ext, 0) + 1

            self.message_queue.put(("log", f"Run ID: {run_id}"))
            self.message_queue.put(("log", f"Found {total_files} qualifying file(s)."))
            self.message_queue.put(("log", "Machine stats captured for formal PDF report."))

            for ext in sorted(per_type_counts):
                self.message_queue.put(("log", f"Discovered by type | {ext} = {per_type_counts[ext]}"))

            if total_files == 0:
                self.message_queue.put(("log", "No qualifying files found."))
                self.message_queue.put(("done_no_files", None))
                return

            self.message_queue.put(("progress_max", total_files))

            rows: List[Dict[str, int | str]] = []
            prefix_summary: Dict[str, int] = {}
            ok_count = 0
            failed_count = 0

            for index, file_path in enumerate(matching_files, start=1):
                file_name = file_path.name
                ext = file_path.suffix.lower()
                prefix = extract_prefix(file_name)

                self.message_queue.put(("log", f"Processing file {index} of {total_files}: {file_name}"))

                try:
                    extractor = get_extractor_for_extension(ext)
                    text = extractor(file_path)

                    row: Dict[str, int | str] = {
                        "File Name": file_name,
                        "Prefix": prefix,
                        "File Type": ext,
                    }

                    hit_line_parts = []
                    for keyword, pattern in patterns.items():
                        count = len(pattern.findall(text))
                        row[keyword] = count
                        hit_line_parts.append(f"{keyword}={count}")

                    row["Total Hits"] = compute_total_hits(row, keywords)
                    rows.append(row)
                    ok_count += 1

                    prefix_summary[prefix] = prefix_summary.get(prefix, 0) + int(row["Total Hits"])

                    if int(row["Total Hits"]) > 0:
                        self.message_queue.put(
                            (
                                "log",
                                f"Hits | {file_name} | Prefix={prefix} | Type={ext} | "
                                + " | ".join(hit_line_parts)
                                + f" | Total={row['Total Hits']}",
                            )
                        )
                    else:
                        self.message_queue.put(
                            ("log", f"Hits | {file_name} | Prefix={prefix} | Type={ext} | Total=0")
                        )

                except Exception as exc:
                    failed_count += 1
                    row = {
                        "File Name": file_name,
                        "Prefix": prefix,
                        "File Type": ext,
                    }
                    for keyword in keywords:
                        row[keyword] = 0
                    row["Total Hits"] = 0
                    rows.append(row)

                    self.message_queue.put(
                        (
                            "log",
                            f"Warning: could not read {file_name} | Type={ext} | {exc}",
                        )
                    )

                self.message_queue.put(("progress_value", index))

            rows = sort_rows_by_prefix_and_total(rows)
            total_keyword_hits = sum(int(row["Total Hits"]) for row in rows)

            self.message_queue.put(("log", "Writing CSV output..."))
            write_csv(output_csv, keywords, rows)
            self.message_queue.put(("log", f"CSV saved: {output_csv}"))

            self.message_queue.put(("log", "Prefix total summary:"))
            for prefix in sorted(prefix_summary):
                self.message_queue.put(("log", f"  {prefix}: {prefix_summary[prefix]} total keyword hit(s)"))

            end_dt = datetime.now()
            elapsed_s = time.perf_counter() - start_perf
            docs_per_min = (total_files / elapsed_s * 60.0) if elapsed_s > 0 else 0.0
            hits_stats = build_hits_stats(rows)

            stats = RunStats(
                run_id=run_id,
                timestamp_compact=timestamp_compact,
                timestamp_human=timestamp_human,
                start_dt=start_dt,
                end_dt=end_dt,
                elapsed_s=elapsed_s,
                total_discovered=total_files,
                processed=len(rows),
                ok=ok_count,
                failed=failed_count,
                docs_per_min=docs_per_min,
                total_keyword_hits=total_keyword_hits,
                keywords=keywords,
                file_types_selected=selected_types,
                input_folder=str(paths.input_folder),
                output_folder=str(paths.output_folder),
                host_name=socket.gethostname(),
                user_name=getpass.getuser(),
                operating_system=platform.platform(),
                python_version=sys.version.split()[0],
                processor_count=os.cpu_count() or 0,
                machine_arch=platform.machine(),
                per_type_counts=per_type_counts,
                prefix_summary=prefix_summary,
                hits_stats=hits_stats,
            )

            self.message_queue.put(("log", "Run Statistics:"))
            self.message_queue.put(("log", f"  host_name={stats.host_name}"))
            self.message_queue.put(("log", f"  user_name={stats.user_name}"))
            self.message_queue.put(("log", f"  operating_system={stats.operating_system}"))
            self.message_queue.put(("log", f"  python_version={stats.python_version}"))
            self.message_queue.put(("log", f"  cpu_count={stats.processor_count}"))
            self.message_queue.put(("log", f"  machine_arch={stats.machine_arch}"))
            self.message_queue.put(("log", f"  start_time={stats.start_dt.strftime('%Y-%m-%d %H:%M:%S')}"))
            self.message_queue.put(("log", f"  end_time={stats.end_dt.strftime('%Y-%m-%d %H:%M:%S')}"))
            self.message_queue.put(("log", f"  elapsed_s={safe_float(stats.elapsed_s)}"))
            self.message_queue.put(("log", f"  total_discovered={stats.total_discovered}"))
            self.message_queue.put(("log", f"  processed={stats.processed}"))
            self.message_queue.put(("log", f"  ok={stats.ok}"))
            self.message_queue.put(("log", f"  failed={stats.failed}"))
            self.message_queue.put(("log", f"  files_per_min={safe_float(stats.docs_per_min)}"))
            self.message_queue.put(("log", f"  total_keyword_hits={stats.total_keyword_hits}"))
            self.message_queue.put(
                (
                    "log",
                    "  total_hits_stats="
                    f"min={safe_float(stats.hits_stats['min'])}, "
                    f"median={safe_float(stats.hits_stats['median'])}, "
                    f"mean={safe_float(stats.hits_stats['mean'])}, "
                    f"max={safe_float(stats.hits_stats['max'])}",
                )
            )

            self.message_queue.put(("log", "Writing PDF event log/report..."))

            # Use all log lines accumulated so far plus the final lines below.
            final_log_lines = self.log_lines.copy()
            final_log_lines.append(timestamped_log_message(f"CSV saved: {output_csv}"))
            final_log_lines.append(timestamped_log_message("Prefix total summary:"))
            for prefix in sorted(prefix_summary):
                final_log_lines.append(
                    timestamped_log_message(f"  {prefix}: {prefix_summary[prefix]} total keyword hit(s)")
                )
            final_log_lines.append(timestamped_log_message("Run Statistics:"))
            final_log_lines.append(timestamped_log_message(f"  host_name={stats.host_name}"))
            final_log_lines.append(timestamped_log_message(f"  user_name={stats.user_name}"))
            final_log_lines.append(timestamped_log_message(f"  operating_system={stats.operating_system}"))
            final_log_lines.append(timestamped_log_message(f"  python_version={stats.python_version}"))
            final_log_lines.append(timestamped_log_message(f"  cpu_count={stats.processor_count}"))
            final_log_lines.append(timestamped_log_message(f"  machine_arch={stats.machine_arch}"))
            final_log_lines.append(
                timestamped_log_message(f"  start_time={stats.start_dt.strftime('%Y-%m-%d %H:%M:%S')}")
            )
            final_log_lines.append(
                timestamped_log_message(f"  end_time={stats.end_dt.strftime('%Y-%m-%d %H:%M:%S')}")
            )
            final_log_lines.append(timestamped_log_message(f"  elapsed_s={safe_float(stats.elapsed_s)}"))
            final_log_lines.append(timestamped_log_message(f"  total_discovered={stats.total_discovered}"))
            final_log_lines.append(timestamped_log_message(f"  processed={stats.processed}"))
            final_log_lines.append(timestamped_log_message(f"  ok={stats.ok}"))
            final_log_lines.append(timestamped_log_message(f"  failed={stats.failed}"))
            final_log_lines.append(timestamped_log_message(f"  files_per_min={safe_float(stats.docs_per_min)}"))
            final_log_lines.append(timestamped_log_message(f"  total_keyword_hits={stats.total_keyword_hits}"))
            final_log_lines.append(
                timestamped_log_message(
                    "  total_hits_stats="
                    f"min={safe_float(stats.hits_stats['min'])}, "
                    f"median={safe_float(stats.hits_stats['median'])}, "
                    f"mean={safe_float(stats.hits_stats['mean'])}, "
                    f"max={safe_float(stats.hits_stats['max'])}"
                )
            )
            final_log_lines.append(timestamped_log_message("Writing PDF event log/report..."))

            build_pdf_report(output_pdf, stats, final_log_lines)

            self.message_queue.put(("log", f"Done. PDF log saved: {output_pdf}"))
            self.message_queue.put(("log", "Run Summary:"))
            self.message_queue.put(
                (
                    "log",
                    f"Completed keyword search across {stats.total_discovered} discovered file(s) "
                    f"in {safe_float(stats.elapsed_s)} seconds.",
                )
            )
            self.message_queue.put(
                (
                    "log",
                    f"A total of {stats.total_keyword_hits} keyword hit(s) were identified across "
                    f"{stats.ok} successfully processed file(s).",
                )
            )
            self.message_queue.put(
                (
                    "log",
                    "Output CSV was grouped by file prefix and sorted within each prefix by highest Total Hits.",
                )
            )
            self.message_queue.put(
                ("log", f"{stats.failed} file(s) could not be read and were recorded as warnings.")
            )
            self.message_queue.put(("done_success", None))

        except Exception as exc:
            self.message_queue.put(("log", f"Fatal error: {exc}"))
            self.message_queue.put(("log", traceback.format_exc()))
            self.message_queue.put(("done_error", None))

    def _poll_queue(self) -> None:
        try:
            while True:
                item_type, payload = self.message_queue.get_nowait()

                if item_type == "log":
                    self._log(str(payload))
                elif item_type == "progress_max":
                    self.progress_bar["maximum"] = int(payload)
                elif item_type == "progress_value":
                    self.progress_bar["value"] = int(payload)
                    current = int(payload)
                    maximum = int(self.progress_bar["maximum"])
                    self.progress_label.configure(text=f"Running ({current}/{maximum})")
                elif item_type == "done_success":
                    self._set_running_state(False)
                    self.progress_label.configure(text="Done")
                    self.progress_bar["value"] = self.progress_bar["maximum"]
                elif item_type == "done_no_files":
                    self._set_running_state(False)
                    self.progress_label.configure(text="No files found")
                    self.progress_bar["value"] = 0
                elif item_type == "done_error":
                    self._set_running_state(False)
                    self.progress_label.configure(text="Error")
        except queue.Empty:
            pass

        self.root.after(100, self._poll_queue)

    def _log(self, message: str) -> None:
        line = timestamped_log_message(message)
        self.log_lines.append(line)

        self.log_widget.configure(state=tk.NORMAL)
        self.log_widget.insert(tk.END, line + "\n")
        self.log_widget.see(tk.END)
        self.log_widget.configure(state=tk.DISABLED)


def main() -> None:
    root = tk.Tk()
    style = ttk.Style()
    try:
        style.theme_use("clam")
    except tk.TclError:
        pass

    DocumentKeywordSearchApp(root)
    root.mainloop()


if __name__ == "__main__":
    main()